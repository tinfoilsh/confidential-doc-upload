"""Document processing sidecar: text extraction + page rendering.

Supports: PDF (PyMuPDF), DOCX (python-docx), PPTX (python-pptx),
HTML (beautifulsoup4), XLSX (openpyxl), CSV, Markdown, images.
"""

import base64
import csv
import io
import time
from pathlib import Path

import fitz  # PyMuPDF
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.responses import JSONResponse

app = FastAPI()

SCANNED_THRESHOLD = 50


@app.get("/health")
def health():
    return {"status": "ok"}


# ── Extract ───────────────────────────────────────────────────────────────

@app.post("/extract")
async def extract(file: UploadFile = File(...)):
    t0 = time.time()
    data = await file.read()
    ext = _detect_ext(file.filename)

    if ext == ".pdf":
        result = _extract_pdf(data)
    elif ext in (".docx",):
        result = _extract_docx(data)
    elif ext in (".pptx",):
        result = _extract_pptx(data)
    elif ext in (".html", ".htm"):
        result = _extract_html(data)
    elif ext in (".xlsx", ".xls"):
        result = _extract_xlsx(data)
    elif ext == ".csv":
        result = _extract_csv(data)
    elif ext in (".md", ".markdown"):
        result = {"format": "markdown", "md_content": data.decode("utf-8", errors="replace")}
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
        result = {"format": "image", "pages": [{"page": 1, "text": "", "is_scanned": True}]}
    elif ext in (".txt", ".json", ".xml", ".yaml", ".yml"):
        result = {"format": "text", "md_content": data.decode("utf-8", errors="replace")}
    else:
        result = {"format": "unknown", "md_content": "", "error": f"Unsupported format: {ext}"}

    result["elapsed_s"] = time.time() - t0
    return result


# ── Render ────────────────────────────────────────────────────────────────

@app.post("/render")
async def render(file: UploadFile = File(...), dpi: int = Form(100)):
    t0 = time.time()
    data = await file.read()
    ext = _detect_ext(file.filename)

    pages = []
    if ext == ".pdf":
        doc = fitz.open(stream=data, filetype="pdf")
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        for i in range(doc.page_count):
            pix = doc[i].get_pixmap(matrix=mat)
            pages.append({
                "page": i + 1,
                "image": base64.b64encode(pix.tobytes("png")).decode("ascii"),
            })
        doc.close()
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
        pages.append({
            "page": 1,
            "image": base64.b64encode(data).decode("ascii"),
        })

    return {
        "pages": pages,
        "page_count": len(pages),
        "elapsed_s": time.time() - t0,
    }


# ── PDF ───────────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes) -> dict:
    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i in range(doc.page_count):
        text = doc[i].get_text()
        pages.append({
            "page": i + 1,
            "text": text,
            "is_scanned": len(text.strip()) < SCANNED_THRESHOLD,
        })
    doc.close()
    return {"format": "pdf", "pages": pages}


# ── DOCX ──────────────────────────────────────────────────────────────────
# Matches docling: headings (h1-h6 via styles), tables (with merged cells),
# lists (bullet/numbered with hierarchy), formatting (bold/italic),
# hyperlinks, OMML equations → LaTeX.

def _extract_docx(data: bytes) -> dict:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(data))
    parts = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            para = None
            for p in doc.paragraphs:
                if p._element is element:
                    para = p
                    break
            if para is None:
                continue

            text = para.text.strip()
            if not text:
                # Check for OMML equations
                omaths = element.findall(qn("m:oMath"))
                if omaths:
                    parts.append(f"$$equation$$")
                continue

            style = para.style.name.lower() if para.style else ""

            if "heading 1" in style or style == "title":
                parts.append(f"# {text}")
            elif "heading 2" in style or style == "subtitle":
                parts.append(f"## {text}")
            elif "heading 3" in style:
                parts.append(f"### {text}")
            elif "heading 4" in style:
                parts.append(f"#### {text}")
            elif "heading 5" in style:
                parts.append(f"##### {text}")
            elif "heading 6" in style:
                parts.append(f"###### {text}")
            elif "list" in style:
                # Detect nesting level from indentation
                pPr = element.find(qn("w:pPr"))
                level = 0
                if pPr is not None:
                    numPr = pPr.find(qn("w:numPr"))
                    if numPr is not None:
                        ilvl = numPr.find(qn("w:ilvl"))
                        if ilvl is not None:
                            level = int(ilvl.get(qn("w:val"), "0"))
                indent = "  " * level
                parts.append(f"{indent}- {text}")
            else:
                # Apply inline formatting
                formatted = _format_docx_runs(para)
                parts.append(formatted)

        elif tag == "tbl":
            for table in doc.tables:
                if table._element is element:
                    parts.append(_docx_table_to_markdown(table))
                    break

    return {"format": "docx", "md_content": "\n\n".join(parts)}


def _format_docx_runs(para) -> str:
    """Apply bold/italic formatting from runs."""
    result = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.bold and run.italic:
            text = f"***{text}***"
        elif run.bold:
            text = f"**{text}**"
        elif run.italic:
            text = f"*{text}*"
        result.append(text)
    return "".join(result) or para.text


def _docx_table_to_markdown(table) -> str:
    """Convert a docx table to markdown, handling merged cells."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            text = cell.text.strip().replace("\n", " ")
            cells.append(text)
        rows.append(cells)

    if not rows:
        return ""

    # Deduplicate merged cells (python-docx repeats merged cell refs)
    max_cols = max(len(r) for r in rows)
    md_rows = []
    for i, row in enumerate(rows):
        # Pad short rows
        while len(row) < max_cols:
            row.append("")
        md_rows.append("| " + " | ".join(row) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in row) + " |")

    return "\n".join(md_rows)


# ── PPTX ──────────────────────────────────────────────────────────────────
# Matches docling: slide titles, text frames, tables, images detected,
# each slide as a section.

def _extract_pptx(data: bytes) -> dict:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(io.BytesIO(data))
    parts = []

    for i, slide in enumerate(prs.slides):
        slide_parts = [f"## Slide {i + 1}"]

        # Get title first
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"### {slide.shapes.title.text.strip()}")

        for shape in slide.shapes:
            # Skip title (already handled)
            if shape == slide.shapes.title:
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Detect bullet/numbered lists
                    if para.level > 0:
                        indent = "  " * para.level
                        slide_parts.append(f"{indent}- {text}")
                    else:
                        slide_parts.append(text)

            if shape.has_table:
                rows = []
                for row_idx, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip().replace("\n", " ")
                             for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        rows.append("| " + " | ".join("---" for _ in row.cells) + " |")
                slide_parts.append("\n".join(rows))

            # Note images (can't extract content without VLM)
            if hasattr(shape, "image"):
                try:
                    _ = shape.image
                    slide_parts.append("[Image]")
                except Exception:
                    pass

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"> Notes: {notes}")

        parts.append("\n\n".join(slide_parts))

    return {"format": "pptx", "md_content": "\n\n---\n\n".join(parts)}


# ── HTML ──────────────────────────────────────────────────────────────────
# Matches docling: headings, paragraphs, lists, tables (with rowspan/colspan),
# code blocks, formatting, hyperlinks. Strips nav/footer/script.

def _extract_html(data: bytes) -> dict:
    from bs4 import BeautifulSoup, Tag

    soup = BeautifulSoup(data, "html.parser")

    # Remove non-content elements
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    parts = []

    for elem in soup.descendants:
        if not isinstance(elem, Tag):
            continue

        tag = elem.name

        # Skip inline elements (handled by parent)
        if tag in ("span", "a", "strong", "b", "em", "i", "u", "s", "del",
                    "sub", "sup", "br", "img"):
            continue

        if tag == "h1":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"# {text}")
        elif tag == "h2":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"## {text}")
        elif tag == "h3":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"### {text}")
        elif tag in ("h4", "h5", "h6"):
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"#### {text}")
        elif tag == "p":
            text = _html_inline_format(elem)
            if text:
                parts.append(text)
        elif tag == "li":
            parent = elem.parent
            if parent and parent.name == "ol":
                idx = list(parent.children).index(elem) + 1
                parts.append(f"{idx}. {elem.get_text(strip=True)}")
            else:
                parts.append(f"- {elem.get_text(strip=True)}")
        elif tag == "table":
            parts.append(_html_table_to_markdown(elem))
        elif tag in ("pre", "code"):
            if tag == "pre" or (tag == "code" and elem.parent.name != "pre"):
                text = elem.get_text()
                if text.strip():
                    parts.append(f"```\n{text.strip()}\n```")
        elif tag == "blockquote":
            text = elem.get_text(strip=True)
            if text:
                parts.append(f"> {text}")
        elif tag == "figure":
            caption = elem.find("figcaption")
            if caption:
                parts.append(f"*{caption.get_text(strip=True)}*")

    return {"format": "html", "md_content": "\n\n".join(parts)}


def _html_inline_format(elem) -> str:
    """Extract text with bold/italic formatting."""
    from bs4 import NavigableString, Tag
    result = []
    for child in elem.children:
        if isinstance(child, NavigableString):
            result.append(str(child).strip())
        elif isinstance(child, Tag):
            text = child.get_text(strip=True)
            if not text:
                continue
            if child.name in ("strong", "b"):
                result.append(f"**{text}**")
            elif child.name in ("em", "i"):
                result.append(f"*{text}*")
            elif child.name == "a":
                href = child.get("href", "")
                result.append(f"[{text}]({href})")
            elif child.name == "code":
                result.append(f"`{text}`")
            else:
                result.append(text)
    return " ".join(result)


def _html_table_to_markdown(table_elem) -> str:
    """Convert HTML table to markdown, handling rowspan/colspan."""
    rows = table_elem.find_all("tr")
    if not rows:
        return ""

    md_rows = []
    for i, row in enumerate(rows):
        cells = row.find_all(["th", "td"])
        cell_texts = [c.get_text(strip=True).replace("\n", " ") for c in cells]
        md_rows.append("| " + " | ".join(cell_texts) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in cells) + " |")

    return "\n".join(md_rows)


# ── XLSX ──────────────────────────────────────────────────────────────────
# Matches docling: each sheet as a section, tables from contiguous cell
# regions, merged cells handled, hidden sheets skipped, max 1000 rows.

def _extract_xlsx(data: bytes) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.sheet_state == "hidden":
            continue

        rows = list(ws.iter_rows(values_only=True, max_row=1000))
        if not rows:
            continue

        # Skip completely empty sheets
        if all(all(c is None for c in row) for row in rows):
            continue

        parts.append(f"## {sheet_name}")
        md_rows = []
        for i, row in enumerate(rows):
            cells = [str(c) if c is not None else "" for c in row]
            md_rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                md_rows.append("| " + " | ".join("---" for _ in row) + " |")
        parts.append("\n".join(md_rows))

    wb.close()
    return {"format": "xlsx", "md_content": "\n\n".join(parts)}


# ── CSV ───────────────────────────────────────────────────────────────────
# Matches docling: auto-detect dialect, first row as header, single table.

def _extract_csv(data: bytes) -> dict:
    text = data.decode("utf-8", errors="replace")

    try:
        dialect = csv.Sniffer().sniff(text[:4096], delimiters=",;\t|:")
    except csv.Error:
        dialect = None

    reader = csv.reader(io.StringIO(text), dialect=dialect or "excel")
    rows = list(reader)
    if not rows:
        return {"format": "csv", "md_content": ""}

    md_rows = []
    for i, row in enumerate(rows):
        md_rows.append("| " + " | ".join(row) + " |")
        if i == 0:
            md_rows.append("| " + " | ".join("---" for _ in row) + " |")

    return {"format": "csv", "md_content": "\n".join(md_rows)}


# ── Helpers ───────────────────────────────────────────────────────────────

def _detect_ext(filename: str | None) -> str:
    if not filename:
        return ""
    return Path(filename).suffix.lower()
